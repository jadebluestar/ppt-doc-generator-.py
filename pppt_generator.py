from pptx import Presentation
from docx import Document

# Replace this list with all your 50 questions
qa_data = [
   
    ("Which company developed the first commercially available quantum computer?",
     ["A) IBM", "B) Google", "C) D-Wave", "D) Microsoft"], "C) D-Wave (2011)"),

    ("India's semiconductor mission targets chip production by which year?",
     ["A) 2025", "B) 2026", "C) 2027", "D) 2028"], "B) 2026"),

    ("The \"Jevons Paradox\" in engineering states that:",
     ["A) Efficiency improvements increase consumption", "B) More engineers reduce innovation", "C) Complexity always fails", "D) Simpler is always better"], "A) Efficiency improvements increase consumption"),

    ("OpenAI's GPT-4 has approximately how many parameters?",
     ["A) 175 billion", "B) 500 billion", "C) 1 trillion", "D) Not publicly disclosed"], "D) Not publicly disclosed"),

    ("In thermodynamics, which is FALSE about entropy?",
     ["A) Always increases in isolated systems", "B) Can decrease locally", "C) Is reversible in all processes", "D) Measures disorder"], "C) Is reversible in all processes"),

    ("Tesla's \"Full Self Driving\" is actually:",
     ["A) Level 5 autonomy", "B) Level 4 autonomy", "C) Level 3 autonomy", "D) Level 2 autonomy"], "D) Level 2 autonomy (still requires driver attention)"),

    ("India's UPI processes approximately how many transactions monthly (2024)?",
     ["A) 5 billion", "B) 10 billion", "C) 15 billion", "D) 20 billion"], "C) ~15 billion"),

    ("Which material has NEGATIVE Poisson's ratio?",
     ["A) Steel", "B) Rubber", "C) Auxetic materials", "D) Copper"], "C) Auxetic materials"),

    ("ChatGPT was released to public in which month/year?",
     ["A) Nov 2022", "B) Dec 2022", "C) Jan 2023", "D) March 2023"], "A) November 2022"),

    ("The \"cold welding\" phenomenon occurs in:",
     ["A) Deep ocean", "B) Antarctica", "C) Outer space", "D) Superconductors"], "C) Outer space (vacuum causes metal atoms to bond)"),

    ("Which country banned TikTok nationwide in 2024?",
     ["A) USA", "B) India", "C) Australia", "D) Canada"], "A) USA (signed into law in 2024)"),

    ("Gaganyaan mission aims to send how many astronauts initially?",
     ["A) 1", "B) 2", "C) 3", "D) 4"], "C) 3 astronauts"),

    ("In aerodynamics, \"coffin corner\" refers to:",
     ["A) Crash zone", "B) Where stall speed meets maximum speed", "C) Landing difficulty", "D) Engine failure zone"], "B) Where stall speed equals maximum speed"),

    ("The \"halting problem\" in computer science is:",
     ["A) Solvable", "B) Unsolvable", "C) Partially solvable", "D) Only solvable with AI"], "B) Unsolvable (proven by Turing)"),

    ("SpaceX's Starship uses which fuel combination?",
     ["A) RP-1/LOX", "B) Hydrogen/LOX", "C) Methane/LOX", "D) Hypergolic"], "C) Methane/LOX"),

    ("India's first 5G testbed was established in which city?",
     ["A) Bangalore", "B) Hyderabad", "C) Multiple IITs", "D) Delhi"], "C) Multiple IITs (IIT Madras, Delhi, Hyderabad, Bombay, Kanpur)"),

    ("The \"observer effect\" in quantum mechanics means:",
     ["A) Watching changes outcome", "B) Measurement affects state", "C) Observers create reality", "D) Nothing changes"], "B) Measurement affects the state"),

    ("Apple's M-series chips are based on which architecture?",
     ["A) x86", "B) ARM", "C) RISC-V", "D) MIPS"], "B) ARM architecture"),

    ("The \"critical angle\" in fiber optics causes:",
     ["A) Light absorption", "B) Total internal reflection", "C) Refraction", "D) Diffraction"], "B) Total internal reflection"),

    ("India's National Quantum Mission budget is approximately:",
     ["A) ₹3,000 crores", "B) ₹6,000 crores", "C) ₹9,000 crores", "D) ₹12,000 crores"], "B) ₹6,000 crores"),

    ("The \"impossible\" EM Drive claimed to produce thrust without propellant. Current scientific consensus?",
     ["A) It works", "B) Measurement error", "C) Classified technology", "D) Still testing"], "B) Measurement error (doesn't work)"),

    ("Nvidia's H100 GPU uses which manufacturing process?",
     ["A) 3nm", "B) 4nm", "C) 5nm", "D) 7nm"], "B) 4nm (TSMC)"),

    ("The \"skin effect\" in conductors means:",
     ["A) Surface corrosion", "B) AC current flows near surface", "C) Heat dissipation", "D) Insulation breakdown"], "B) AC current concentrates near surface"),

    ("India's Aditya-L1 mission studies the Sun from which Lagrange point?",
     ["A) L1", "B) L2", "C) L3", "D) L4"], "A) L1 (between Earth and Sun)"),

    ("The \"dark silicon\" problem in processors means:",
     ["A) Light-sensitive chips", "B) Can't power all transistors simultaneously", "C) Black silicon wafers", "D) Undetectable backdoors"], "B) Can't power all transistors due to heat"),

    ("CERN's Large Hadron Collider discovered the Higgs Boson in which year?",
     ["A) 2010", "B) 2011", "C) 2012", "D) 2013"], "C) 2012"),

    ("The \"tyranny of the rocket equation\" means:",
     ["A) Rockets are expensive", "B) Exponential mass increase for linear velocity", "C) Physics prevents space travel", "D) Fuel efficiency limits"], "B) Exponential mass increase for linear delta-v"),

    ("Meta's Llama 3 model is:",
     ["A) Closed source", "B) Open source", "C) Partially open", "D) Government classified"], "B) Open source"),

    ("The \"Kármán line\" (space boundary) is at what altitude?",
     ["A) 50 km", "B) 80 km", "C) 100 km", "D) 150 km"], "C) 100 km"),

    ("India's semiconductor fabrication plants are being set up in which states primarily?",
     ["A) Karnataka, TN", "B) Gujarat, Assam", "C) Maharashtra, UP", "D) All of these"], "B) Gujarat, Assam (major announcements)"),

    ("The \"von Neumann bottleneck\" refers to:",
     ["A) CPU speed", "B) Memory-CPU data transfer limitation", "C) Power consumption", "D) Heat dissipation"], "B) Memory-CPU data transfer bottleneck"),

    ("TSMC's most advanced production node (as of 2024) is:",
     ["A) 2nm", "B) 3nm", "C) 4nm", "D) 5nm"], "B) 3nm (N3)"),

    ("The \"bootstrap paradox\" in engineering relates to:",
     ["A) Computer booting", "B) Self-causing loop", "C) Startup problems", "D) Power-on sequence"], "B) Self-causing causal loop"),

    ("Anthropic's Claude uses which safety technique primarily?",
     ["A) RLHF", "B) Constitutional AI", "C) Supervised learning", "D) Rule-based"], "B) Constitutional AI"),

    ("The \"Dunning-Kruger effect\" is most relevant in:",
     ["A) Material failure", "B) Overconfidence in low-skill individuals", "C) Structural engineering", "D) Circuit design"], "B) Overconfidence in incompetent individuals"),

    ("RISC-V is significant because it's:",
     ["A) Fastest architecture", "B) Open-source ISA", "C) Chinese technology", "D) Low power only"], "B) Open-source instruction set"),

    ("The \"AI alignment problem\" refers to:",
     ["A) GPU synchronization", "B) Aligning AI goals with human values", "C) Neural network training", "D) Data labeling"], "B) Aligning AI goals with human values"),

    ("India's first private rocket launch was by which company?",
     ["A) Skyroot", "B) Agnikul", "C) Bellatrix", "D) Pixxel"], "A) Skyroot Aerospace (Vikram-S, 2022)"),

    ("The \"p vs np\" problem in computer science is:",
     ["A) Solved", "B) Proven impossible", "C) Still unsolved", "D) Partially solved"], "C) Still unsolved (Clay Millennium Prize)"),

    ("Neuralink's brain implant has how many electrodes approximately?",
     ["A) 100", "B) 1,000", "C) 10,000", "D) 100,000"], "B) ~1,000 electrodes"),

    ("The \"Casimir effect\" demonstrates:",
     ["A) Vacuum energy", "B) Quantum tunneling", "C) Wave-particle duality", "D) Entanglement"], "A) Vacuum energy creates force"),

    ("Google's latest AI model Gemini was released in which year?",
     ["A) 2022", "B) 2023", "C) 2024", "D) Not yet released"], "B) December 2023"),

    ("The \"deadlock\" condition in operating systems requires how many conditions?",
     ["A) 2", "B) 3", "C) 4", "D) 5"], "C) 4 conditions (Mutual exclusion, Hold and wait, No preemption, Circular wait)"),

    ("India's fastest train Vande Bharat runs at max operational speed of:",
     ["A) 130 km/h", "B) 160 km/h", "C) 180 km/h", "D) 200 km/h"], "B) 160 km/h (designed for 180 km/h)"),

    ("The \"Shannon limit\" defines maximum:",
     ["A) Internet speed", "B) Channel capacity", "C) Processor frequency", "D) Storage density"], "B) Channel capacity for given noise"),

    ("Waabi's self-driving approach uses primarily:",
     ["A) Cameras only", "B) LiDAR only", "C) AI simulation", "D) HD maps"], "C) AI simulation-first approach"),

    ("The \"birthday paradox\" shows that in a group of 23 people, probability of shared birthday is:",
     ["A) 23%", "B) ~50%", "C) 75%", "D) 90%"], "B) ~50%"),

    ("India's National Research Foundation budget is approximately:",
     ["A) ₹10,000 crores", "B) ₹25,000 crores", "C) ₹50,000 crores", "D) ₹1,00,000 crores"], "C) ₹50,000 crores over 5 years"),

    ("The \"Byzantine Generals Problem\" is solved by:",
     ["A) Encryption", "B) Blockchain consensus", "C) Firewalls", "D) AI"], "B) Blockchain consensus mechanisms"),

    ("Sam Altman was briefly fired from OpenAI in which month/year?",
     ["A) Oct 2023", "B) Nov 2023", "C) Dec 2023", "D) Jan 2024"], "B) November 2023 (reinstated days later)")
] 


# --- PPT with options ---
ppt_with_options = Presentation()
for i, (q, opts, ans) in enumerate(qa_data, 1):
    slide = ppt_with_options.slides.add_slide(ppt_with_options.slide_layouts[1])
    slide.shapes.title.text = f"Q{i}"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = q
    for opt in opts:
        tf.add_paragraph().text = opt
ppt_with_options.save("quiz_with_options.pptx")

# --- PPT without options ---
ppt_without_options = Presentation()
for i, (q, opts, ans) in enumerate(qa_data, 1):
    slide = ppt_without_options.slides.add_slide(ppt_without_options.slide_layouts[1])
    slide.shapes.title.text = f"Q{i}"
    slide.shapes.placeholders[1].text = q
ppt_without_options.save("quiz_without_options.pptx")

# --- Word document with Q + A ---
doc = Document()
for i, (q, opts, ans) in enumerate(qa_data, 1):
    doc.add_heading(f"Q{i}. {q}", level=2)
    for opt in opts:
        doc.add_paragraph(opt, style='List Bullet')
    doc.add_paragraph(f"Answer: {ans}", style='Intense Quote')
doc.save("quiz_QA.docx") 